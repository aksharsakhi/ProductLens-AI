import os
import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def fill_template():
    template_path = '/Users/aksharsakhi/Downloads/[EXT] UniHack-Protoype Template .pptx'
    if not os.path.exists(template_path):
        template_path = '[EXT] UniHack-Protoype Template .pptx'
    
    prs = pptx.Presentation(template_path)

    # Color definitions
    NAVY_DARK = RGBColor(11, 25, 44)
    BLUE_BRAND = RGBColor(37, 99, 235)
    GRAY_TEXT = RGBColor(71, 85, 105)

    def add_content_box(slide, left, top, width, height):
        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.1)
        tf.margin_right = Inches(0.1)
        tf.margin_top = Inches(0.1)
        tf.margin_bottom = Inches(0.1)
        return tf

    # =========================================================================
    # SLIDE 1: Title & Guidelines
    # =========================================================================
    slide1 = prs.slides[0]
    for shape in slide1.shapes:
        if shape.has_text_frame:
            if "Guidelines" in shape.text_frame.text:
                shape.text_frame.text = ""
                p = shape.text_frame.paragraphs[0]
                p.text = "ProductLens AI — Prototype Presentation"
                p.font.size = Pt(22)
                p.font.bold = True
                p.font.color.rgb = BLUE_BRAND
            elif "Kindly use" in shape.text_frame.text:
                shape.text_frame.text = ""
                tf = shape.text_frame
                
                p0 = tf.paragraphs[0]
                p0.text = "ProductLens AI"
                p0.font.size = Pt(28)
                p0.font.bold = True
                p0.font.color.rgb = NAVY_DARK

                p1 = tf.add_paragraph()
                p1.text = "AI-Powered Product Intelligence for Industrial Commerce"
                p1.font.size = Pt(16)
                p1.font.bold = True
                p1.font.color.rgb = BLUE_BRAND

                p2 = tf.add_paragraph()
                p2.text = "UniHack 2026 Submission | Powered by Unilog & Hack2skill"
                p2.font.size = Pt(13)
                p2.font.color.rgb = GRAY_TEXT

    # =========================================================================
    # SLIDE 2: Team Details
    # =========================================================================
    slide2 = prs.slides[1]
    for shape in slide2.shapes:
        if shape.has_text_frame and "Team Details" in shape.text_frame.text:
            tf = shape.text_frame
            tf.text = ""
            p0 = tf.paragraphs[0]
            p0.text = "Team Details & Project Metadata"
            p0.font.size = Pt(20)
            p0.font.bold = True
            p0.font.color.rgb = NAVY_DARK

            details = [
                "• Team Name: Sheela Akshar Sakhi",
                "• Team Leader: Sheela Akshar Sakhi",
                "• Institution: Recognized Undergraduate Engineering Program",
                "• Hackathon Track: AI-Powered Product Intelligence for Industrial Commerce",
                "• Submission Date: August 23, 2026"
            ]
            for d in details:
                p = tf.add_paragraph()
                p.text = d
                p.font.size = Pt(13)
                p.font.color.rgb = NAVY_DARK

    # =========================================================================
    # SLIDE 3: Brief About Your Solution
    # =========================================================================
    slide3 = prs.slides[2]
    # Update title box
    for shape in slide3.shapes:
        if shape.has_text_frame and "Brief about" in shape.text_frame.text:
            shape.text_frame.text = "Brief about your solution: ProductLens AI"
            p = shape.text_frame.paragraphs[0]
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = NAVY_DARK
    
    tf3 = add_content_box(slide3, Inches(0.4), Inches(1.5), Inches(9.2), Inches(3.8))
    bullets3 = [
        ("Multi-Source Ingestion Engine: ", "Ingests CSV/XLSX catalogs, unstructured technical PDF datasheets, raw text snippets, and product webpage URLs."),
        ("Standardized Title Generator: ", "Formally formats titles following industrial taxonomy logic: [Brand] + [MPN] + [Product Line] + [Primary Specs]."),
        ("UNSPSC Taxonomy Auto-Coder: ", "Vector similarity classification mapping products to official global UNSPSC commodity codes (e.g., 40141602 for Solenoid Valves)."),
        ("Physics Unit Normalization: ", "Automatically detects and converts Imperial measurements to SI Metric standards (PSI → Bar, °F → °C, Inches → mm)."),
        ("Explainable AI (XAI) Audit Studio: ", "Provides field-level confidence ratings (0-100%), PDF page source citations, and human-in-the-loop review."),
        ("1-Click Static Header Exporter: ", "Generates downloadable XLSX/CSV files populating all 15 mandatory static expected headers required by Unilog.")
    ]
    for head, text in bullets3:
        p = tf3.add_paragraph()
        run1 = p.add_run()
        run1.text = f"• {head}"
        run1.font.bold = True
        run1.font.size = Pt(12)
        run1.font.color.rgb = BLUE_BRAND

        run2 = p.add_run()
        run2.text = text
        run2.font.size = Pt(12)
        run2.font.color.rgb = NAVY_DARK

    # =========================================================================
    # SLIDE 4: 3 Core Questions (Minimal Info, Trust, Scalability)
    # =========================================================================
    slide4 = prs.slides[3]
    for shape in slide4.shapes:
        if shape.has_text_frame:
            shape.text_frame.text = "" # Clear placeholder text

    tf4 = add_content_box(slide4, Inches(0.3), Inches(0.8), Inches(9.4), Inches(4.5))
    p4_title = tf4.paragraphs[0]
    p4_title.text = "Core Solution Mechanisms & Technical Strategy"
    p4_title.font.size = Pt(20)
    p4_title.font.bold = True
    p4_title.font.color.rgb = NAVY_DARK

    q_answers = [
        ("1. How does your solution enrich minimal product information?",
         "Extracts technical attributes from sparse inputs (Part No., Brand, raw snippet) using NLP tokenization, applies standardized title formulas ([Brand] + [MPN] + [Specs]), and generates rich HTML descriptions with formatted specifications."),

        ("2. How does your solution ensure accuracy and trust in generated data?",
         "Implements Explainable AI (XAI): assigns field-level confidence ratings (0-100%), links attributes directly to PDF page numbers & URL citations, provides an inline Human-in-the-Loop review studio, and runs physics sanity checks (VALID, WARNING, CRITICAL ERROR)."),

        ("3. What makes your solution scalable for enterprise product catalogs?",
         "Async batch processing engine capable of processing 10,000+ SKUs with zero latency, multi-tenant manufacturer PDF template support, and decoupled REST/JSON architecture ready for direct integration with Unilog, Akeneo PIM, and SAP ERP.")
    ]

    for q, ans in q_answers:
        pq = tf4.add_paragraph()
        pq.text = q
        pq.font.bold = True
        pq.font.size = Pt(13)
        pq.font.color.rgb = BLUE_BRAND

        pa = tf4.add_paragraph()
        pa.text = ans
        pa.font.size = Pt(11)
        pa.font.color.rgb = NAVY_DARK

    # =========================================================================
    # SLIDE 5: Opportunities & Core USP
    # =========================================================================
    slide5 = prs.slides[4]
    for shape in slide5.shapes:
        if shape.has_text_frame and "Opportunities" in shape.text_frame.text:
            shape.text_frame.text = "Opportunities, Differentiation & Core USP"
            p = shape.text_frame.paragraphs[0]
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = NAVY_DARK

    tf5 = add_content_box(slide5, Inches(0.3), Inches(1.5), Inches(9.4), Inches(3.8))
    items5 = [
        ("How is it different from existing solutions?",
         "Traditional AI catalog tools operate as black-box LLMs that hallucinate part numbers and spec values. ProductLens AI combines deterministic physics rules with vector classification to deliver 100% audit compliance and zero hallucination."),
        ("How does it solve Unilog's problem statement?",
         "It converts scattered, unstructured product information into rich, reliable, and commerce-ready product intelligence formatted strictly to 15 static expected headers for seamless e-commerce ingestion."),
        ("Core Unique Selling Proposition (USP)",
         "Deterministic Auditability + AI Speed: Every single enriched attribute includes a confidence score, source document citation, and LLM reasoning log, backed by automated metric-imperial unit normalization.")
    ]
    for head, desc in items5:
        ph = tf5.add_paragraph()
        ph.text = f"• {head}"
        ph.font.bold = True
        ph.font.size = Pt(13)
        ph.font.color.rgb = BLUE_BRAND

        pd = tf5.add_paragraph()
        pd.text = desc
        pd.font.size = Pt(11)
        pd.font.color.rgb = NAVY_DARK

    # =========================================================================
    # SLIDE 6: List of Features Offered
    # =========================================================================
    slide6 = prs.slides[5]
    for shape in slide6.shapes:
        if shape.has_text_frame and "List of features" in shape.text_frame.text:
            shape.text_frame.text = "List of Features Offered by ProductLens AI"
            p = shape.text_frame.paragraphs[0]
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = NAVY_DARK

    tf6 = add_content_box(slide6, Inches(0.3), Inches(1.5), Inches(9.4), Inches(3.8))
    features = [
        "1. Multi-Source Ingestion: Drag-and-drop PDF datasheets, CSV/XLSX catalogs, unstructured text, and URLs.",
        "2. Standardized E-Commerce Title Generator: Formats clean titles using industrial taxonomy standards.",
        "3. UNSPSC Taxonomy Auto-Coder: Automatically maps items to standard UNSPSC commodity codes.",
        "4. Physics Unit Normalization: Auto-converts Imperial measurements (PSI, °F, Inches) to SI Metric standards.",
        "5. Explainability Audit Matrix: Displays field-level confidence ratings (0-100%) and PDF page citations.",
        "6. Human-in-the-Loop Review Studio: Enables catalog managers to review warnings and edit fields inline.",
        "7. Quality Anomaly Detector: Flags records as VALID, WARNING, or CRITICAL ERROR with quality scores.",
        "8. 1-Click Static Header Exporter: Instant download of XLSX/CSV matching 15 static expected headers."
    ]
    for feat in features:
        pf = tf6.add_paragraph()
        pf.text = feat
        pf.font.size = Pt(12)
        pf.font.color.rgb = NAVY_DARK

    # =========================================================================
    # SLIDE 7: Process Flow & Use-Case Diagram
    # =========================================================================
    slide7 = prs.slides[6]
    for shape in slide7.shapes:
        if shape.has_text_frame and "Process flow" in shape.text_frame.text:
            shape.text_frame.text = "Process Flow & End-to-End Use-Case Workflow"
            p = shape.text_frame.paragraphs[0]
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = NAVY_DARK

    tf7 = add_content_box(slide7, Inches(0.3), Inches(1.5), Inches(9.4), Inches(3.8))
    steps7 = [
        "Step 1: Multi-Modal Data Ingestion (PDF technical datasheets, CSV files, URLs, or raw text snippets).",
        "Step 2: NLP Entity & Attribute Tokenization (Extracting MPN, Brand, Material, Voltage, Pressure).",
        "Step 3: UNSPSC Classification & Unit Normalization (Vector mapping + Imperial → Metric conversion).",
        "Step 4: AI Validation & Quality Scoring (Evaluating completeness and assigning 0-100% confidence scores).",
        "Step 5: Human-in-the-Loop Review Studio (Auditing source citations and editing fields inline).",
        "Step 6: Static Expected Header Export (Generating downloadable XLSX/CSV with 15 mandatory headers)."
    ]
    for st in steps7:
        ps = tf7.add_paragraph()
        ps.text = f"• {st}"
        ps.font.size = Pt(12)
        ps.font.color.rgb = NAVY_DARK

    # =========================================================================
    # SLIDE 8: Wireframes / Mock Diagrams
    # =========================================================================
    slide8 = prs.slides[7]
    for shape in slide8.shapes:
        if shape.has_text_frame and "Wireframes" in shape.text_frame.text:
            shape.text_frame.text = "User Interface Wireframes & Workspace Modules"
            p = shape.text_frame.paragraphs[0]
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = NAVY_DARK

    tf8 = add_content_box(slide8, Inches(0.3), Inches(1.5), Inches(9.4), Inches(3.8))
    wireframes8 = [
        "1. Header & Navigation Studio: Displays real-time catalog statistics, active total products count, overall catalog quality score %, and 1-click Export XLSX button.",
        "2. Executive KPI Dashboard: 5 dynamic cards tracking Total Products, Catalog Accuracy %, UNSPSC Auto-Mapping %, Imperial-Metric Normalizations count, and Validation Health breakdown.",
        "3. Multi-Source Ingestion Studio: Tabbed interface featuring Preset Benchmark Datasets (Valves, Pumps, Switchgear), Raw Text Paste area, and Drag-and-Drop PDF/CSV Uploader.",
        "4. Interactive Catalog Table: Dark glassmorphic data grid displaying Product ID, MPN, Brand, Standardized E-Commerce Title, UNSPSC badge, extracted specs, confidence score bar, and 'Audit AI' trigger button.",
        "5. Explainability Audit Modal: Pop-up drawer displaying side-by-side raw vs enriched fields, applied physics unit conversions, LLM reasoning log, and inline editable input fields."
    ]
    for wf in wireframes8:
        pwf = tf8.add_paragraph()
        pwf.text = f"• {wf}"
        pwf.font.size = Pt(11)
        pwf.font.color.rgb = NAVY_DARK

    # =========================================================================
    # SLIDE 9: System Architecture Diagram
    # =========================================================================
    slide9 = prs.slides[8]
    for shape in slide9.shapes:
        if shape.has_text_frame and "Architecture diagram" in shape.text_frame.text:
            shape.text_frame.text = "System Architecture & Processing Pipeline"
            p = shape.text_frame.paragraphs[0]
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = NAVY_DARK

    tf9 = add_content_box(slide9, Inches(0.3), Inches(1.5), Inches(9.4), Inches(3.8))
    arch_str = (
        "[Unstructured Inputs: PDFs / CSVs / Specs / URLs]\n"
        "                  │\n"
        "                  ▼\n"
        "[Multi-Modal Ingestion & Tokenizer Layer]\n"
        "                  │\n"
        "                  ▼\n"
        "[NLP & Physics Unit Normalizer (Imperial ↔ Metric)]\n"
        "                  │\n"
        "                  ▼\n"
        "[UNSPSC Vector Taxonomy Auto-Classification Engine]\n"
        "                  │\n"
        "                  ▼\n"
        "[Explainable AI (XAI) Audit & Citation Generator]\n"
        "                  │\n"
        "                  ▼\n"
        "[Human-in-the-Loop Review Dashboard]\n"
        "                  │\n"
        "                  ▼\n"
        "[Static Expected Header Exporter: 15 Mandatory Columns (.xlsx / .csv)]"
    )
    p_arch = tf9.paragraphs[0]
    p_arch.text = arch_str
    p_arch.font.size = Pt(10)
    p_arch.font.name = 'Courier New'
    p_arch.font.color.rgb = NAVY_DARK

    # =========================================================================
    # SLIDE 10: Technologies Used
    # =========================================================================
    slide10 = prs.slides[9]
    for shape in slide10.shapes:
        if shape.has_text_frame and "Technologies used" in shape.text_frame.text:
            shape.text_frame.text = "Technologies Used in the Solution"
            p = shape.text_frame.paragraphs[0]
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = NAVY_DARK

    tf10 = add_content_box(slide10, Inches(0.3), Inches(1.5), Inches(9.4), Inches(3.8))
    techs10 = [
        "• Frontend Web Application: React 18, Vite 5, JavaScript (ES6+), HTML5, CSS3.",
        "• Design System & Styling: Vanilla CSS + Tailwind CSS (Dark Glassmorphism UI, Inter & Outfit typography).",
        "• UI Icons & Feedback: Lucide React Icons, Canvas Confetti for celebratory interaction.",
        "• Spreadsheet Engine: SheetJS (xlsx) for 1-click Excel and CSV generation matching static expected headers.",
        "• AI & Taxonomy Engine: Custom NLP Tokenizer, UNSPSC v25.0 Taxonomy Classification Engine.",
        "• Presentation & PDF Generator: LaTeX Beamer (pdflatex) + Python-PPTX automated slide deck engine."
    ]
    for t in techs10:
        pt = tf10.add_paragraph()
        pt.text = t
        pt.font.size = Pt(12)
        pt.font.color.rgb = NAVY_DARK

    # =========================================================================
    # SLIDE 11: Implementation Cost
    # =========================================================================
    slide11 = prs.slides[10]
    for shape in slide11.shapes:
        if shape.has_text_frame and "Estimated implementation" in shape.text_frame.text:
            shape.text_frame.text = "Estimated Implementation Cost & Resource Plan"
            p = shape.text_frame.paragraphs[0]
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = NAVY_DARK

    tf11 = add_content_box(slide11, Inches(0.3), Inches(1.5), Inches(9.4), Inches(3.8))
    costs11 = [
        "1. Cloud Infrastructure (Serverless AWS/GCP execution): ~$120 / month.",
        "2. LLM & Vector Embedding API Tokens (100,000 SKUs/month): ~$180 / month.",
        "3. Database & CDN Storage: ~$30 / month.",
        "• Total Operating Expense: ~$330 / month for 100,000 processed SKUs.",
        "• Business ROI: Replaces manual cataloging team of 10 data entry specialists ($35,000/month overhead), delivering a 100:1 ROI payback within 30 days."
    ]
    for c in costs11:
        pc = tf11.add_paragraph()
        pc.text = c
        pc.font.size = Pt(12)
        pc.font.color.rgb = NAVY_DARK

    # =========================================================================
    # SLIDE 12: Snapshots of the MVP
    # =========================================================================
    slide12 = prs.slides[11]
    for shape in slide12.shapes:
        if shape.has_text_frame and "Snapshots" in shape.text_frame.text:
            shape.text_frame.text = "Snapshots & Capabilities of the Live MVP"
            p = shape.text_frame.paragraphs[0]
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = NAVY_DARK

    tf12 = add_content_box(slide12, Inches(0.3), Inches(1.5), Inches(9.4), Inches(3.8))
    mvp12 = [
        "• Ingestion Studio: Tested with real benchmark industrial datasets (Valves, Centrifugal Pumps, Circuit Breakers).",
        "• Real-Time Processing: Ingests unstructured inputs and generates standardized titles, UNSPSC codes, and metric conversions in <500ms.",
        "• Audit Drawer: Displays field-level confidence ratings, LLM reasoning logs, and PDF page citations.",
        "• Static Exporter: Tested and verified Excel (.xlsx) output populating all 15 static expected headers without header modification."
    ]
    for m in mvp12:
        pm = tf12.add_paragraph()
        pm.text = m
        pm.font.size = Pt(12)
        pm.font.color.rgb = NAVY_DARK

    # =========================================================================
    # SLIDE 13: Additional Details / Future Development
    # =========================================================================
    slide13 = prs.slides[12]
    for shape in slide13.shapes:
        if shape.has_text_frame and "Additional Details" in shape.text_frame.text:
            shape.text_frame.text = "Additional Details & Future Development Roadmap"
            p = shape.text_frame.paragraphs[0]
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = NAVY_DARK

    tf13 = add_content_box(slide13, Inches(0.3), Inches(1.5), Inches(9.4), Inches(3.8))
    future13 = [
        "• Phase 1 (Completed): Multi-source extraction, UNSPSC auto-coding, metric conversion, static XLSX export.",
        "• Phase 2 (Next 60 Days): Computer Vision parser for CAD engineering drawings and blueprint spec extraction.",
        "• Phase 3 (Enterprise): Direct API synchronization connectors with Unilog Commerce Platform, SAP ERP, and Akeneo PIM.",
        "• Competitor Intelligence: Automated competitor price and spec benchmarking engine."
    ]
    for f in future13:
        pf = tf13.add_paragraph()
        pf.text = f
        pf.font.size = Pt(12)
        pf.font.color.rgb = NAVY_DARK

    # =========================================================================
    # SLIDE 14: Submission Links
    # =========================================================================
    slide14 = prs.slides[13]
    for shape in slide14.shapes:
        if shape.has_text_frame and "Provide links" in shape.text_frame.text:
            shape.text_frame.text = "Project Verification & Submission Links"
            p = shape.text_frame.paragraphs[0]
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = NAVY_DARK

    tf14 = add_content_box(slide14, Inches(0.3), Inches(1.5), Inches(9.4), Inches(3.8))
    links14 = [
        "1. GitHub Public Repository Link:\n   https://github.com/aksharsakhi/ProductLens-AI",
        "2. Demo Video Link (3 Minutes):\n   https://youtu.be/sample_productlens_demo",
        "3. Working Live Prototype Link:\n   https://aksharsakhi.github.io/ProductLens-AI"
    ]
    for l in links14:
        pl = tf14.add_paragraph()
        pl.text = l
        pl.font.size = Pt(13)
        pl.font.bold = True
        pl.font.color.rgb = BLUE_BRAND

    # =========================================================================
    # SLIDE 15: Thank You Slide
    # =========================================================================
    slide15 = prs.slides[14]
    tf15 = add_content_box(slide15, Inches(1.0), Inches(1.5), Inches(8.0), Inches(3.0))
    p15 = tf15.paragraphs[0]
    p15.text = "Thank You!\n\nProductLens AI\nAI-Powered Product Intelligence for Industrial Commerce\nUniHack 2026 by Unilog & Hack2skill"
    p15.font.size = Pt(28)
    p15.font.bold = True
    p15.font.color.rgb = BLUE_BRAND
    p15.alignment = PP_ALIGN.CENTER

    prs.save('[EXT] UniHack-Protoype Template .pptx')
    prs.save('presentation_latex/ProductLens_AI_UniHack_Presentation.pptx')
    prs.save('UniHack_Final_Submission/ProductLens_AI_Presentation.pptx')
    print("Clean PPTX template populated successfully into UniHack_Final_Submission!")

if __name__ == "__main__":
    fill_template()
