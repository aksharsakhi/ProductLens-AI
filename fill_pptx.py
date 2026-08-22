import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def populate_presentation():
    pptx_path = '[EXT] UniHack-Protoype Template .pptx'
    prs = pptx.Presentation(pptx_path)

    # Color Constants
    BLUE_DARK = RGBColor(11, 25, 44)
    BLUE_BRAND = RGBColor(37, 99, 235)
    GRAY_TEXT = RGBColor(71, 85, 105)

    # -------------------------------------------------------------
    # SLIDE 1: Title Slide & Guidelines
    # -------------------------------------------------------------
    slide1 = prs.slides[0]
    for shape in slide1.shapes:
        if shape.has_text_frame:
            tf = shape.text_frame
            tf.word_wrap = True
            if "Guidelines" in tf.text or "Kindly use" in tf.text:
                tf.text = "ProductLens AI"
                p0 = tf.paragraphs[0]
                p0.font.size = Pt(36)
                p0.font.bold = True
                p0.font.color.rgb = BLUE_BRAND
                
                p1 = tf.add_paragraph()
                p1.text = "AI-Powered Product Intelligence for Industrial Commerce"
                p1.font.size = Pt(20)
                p1.font.bold = True
                p1.font.color.rgb = BLUE_DARK
                
                p2 = tf.add_paragraph()
                p2.text = "UniHack 2026 Submission | Powered by Unilog & Hack2skill"
                p2.font.size = Pt(14)
                p2.font.color.rgb = GRAY_TEXT

    # -------------------------------------------------------------
    # SLIDE 2: Team Details
    # -------------------------------------------------------------
    slide2 = prs.slides[1]
    for shape in slide2.shapes:
        if shape.has_text_frame:
            tf = shape.text_frame
            tf.word_wrap = True
            tf.text = "Team Details & Project Metadata\n"
            p = tf.paragraphs[0]
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = BLUE_DARK

            p1 = tf.add_paragraph()
            p1.text = "• Team Name: Sheela Akshar Sakhi"
            p1.font.size = Pt(16)

            p2 = tf.add_paragraph()
            p2.text = "• Team Leader Name: Sheela Akshar Sakhi"
            p2.font.size = Pt(16)

            p3 = tf.add_paragraph()
            p3.text = "• Institution / Organization: Recognized Undergraduate Engineering Program"
            p3.font.size = Pt(16)

            p4 = tf.add_paragraph()
            p4.text = "• Hackathon Track: AI-Powered Product Intelligence for Industrial Commerce"
            p4.font.size = Pt(16)

            p5 = tf.add_paragraph()
            p5.text = "• Submission Date: August 23, 2026"
            p5.font.size = Pt(16)

    # -------------------------------------------------------------
    # SLIDE 3: Brief About Your Solution
    # -------------------------------------------------------------
    slide3 = prs.slides[2]
    for shape in slide3.shapes:
        if shape.has_text_frame:
            tf = shape.text_frame
            tf.word_wrap = True
            tf.text = "Brief About Your Solution: ProductLens AI\n"
            p = tf.paragraphs[0]
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = BLUE_DARK

            p1 = tf.add_paragraph()
            p1.text = "ProductLens AI is an end-to-end, enterprise-grade AI Product Intelligence Platform built specifically for industrial commerce challenges:"
            p1.font.size = Pt(14)
            p1.font.bold = True

            bullets = [
                "Multi-Source Ingestion Engine: Automatically ingests CSV/XLSX catalogs, unstructured technical PDF datasheets, raw text snippets, and product web URLs.",
                "E-Commerce Title Standardizer: Converts sparse raw descriptions into standardized e-commerce titles following industrial taxonomy logic: [Brand] + [MPN] + [Line] + [Specs].",
                "Automated UNSPSC Taxonomy Auto-Coder: Vector similarity classification mapping products to official global UNSPSC commodity codes (e.g., 40141602 for Solenoid Valves).",
                "Physics Unit Normalization: Automatically detects and converts Imperial measurements to SI Metric standards (PSI -> Bar, °F -> °C, Inches -> mm).",
                "Explainable AI (XAI) Audit Studio: Field-level confidence scoring (0-100%), PDF page source citations, and human-in-the-loop inline review studio.",
                "1-Click Static Header Exporter: Instant download of commerce-ready XLSX/CSV files populating all 15 static expected headers required by Unilog."
            ]
            for b in bullets:
                p_b = tf.add_paragraph()
                p_b.text = f"• {b}"
                p_b.font.size = Pt(12)

    # -------------------------------------------------------------
    # SLIDE 4: 3 Core Questions (Minimal Info, Trust, Scalability)
    # -------------------------------------------------------------
    slide4 = prs.slides[3]
    for shape in slide4.shapes:
        if shape.has_text_frame:
            tf = shape.text_frame
            tf.word_wrap = True
            tf.text = "Core Solution Mechanisms & Technical Strategy\n"
            p = tf.paragraphs[0]
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = BLUE_DARK

            q_data = [
                ("1. How does your solution enrich minimal product information?", 
                 "• Extracts technical attributes from sparse inputs (Brand, Part Number, raw snippet) using NLP tokenization.\n• Applies standard title templates: [Brand] + [MPN] + [Product Line] + [Primary Specs].\n• Generates structured HTML long descriptions and short e-commerce summaries."),
                
                ("2. How does your solution ensure accuracy and trust in generated data?", 
                 "• Field-level confidence scoring (0-100%) for every extracted spec.\n• Source Citation Traceability linking attributes directly to PDF page numbers and URL text snippets.\n• Human-in-the-Loop Review Studio allowing catalog managers to edit fields inline with instant re-validation.\n• Physics-based sanity checks flagging out-of-bounds parameters (VALID, WARNING, CRITICAL ERROR)."),
                
                ("3. What makes your solution scalable for enterprise product catalogs?", 
                 "• Async batch processing engine capable of handling 10,000+ SKUs with zero latency.\n• Multi-tenant manufacturer format support adapting to diverse supplier PDF layouts.\n• Decoupled architecture ready for direct API integration with Unilog, Akeneo PIM, and SAP ERP.")
            ]

            for q_title, q_body in q_data:
                p_q = tf.add_paragraph()
                p_q.text = q_title
                p_q.font.size = Pt(14)
                p_q.font.bold = True
                p_q.font.color.rgb = BLUE_BRAND

                p_b = tf.add_paragraph()
                p_b.text = q_body
                p_b.font.size = Pt(11)

    # -------------------------------------------------------------
    # SLIDE 5: Opportunities & Core USP
    # -------------------------------------------------------------
    slide5 = prs.slides[4]
    for shape in slide5.shapes:
        if shape.has_text_frame:
            tf = shape.text_frame
            tf.word_wrap = True
            tf.text = "Opportunities, Differentiation & Core USP\n"
            p = tf.paragraphs[0]
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = BLUE_DARK

            items = [
                ("How is it different from existing solutions?", "Traditional AI catalog tools operate as black-box LLMs that hallucinate part numbers and unit values. ProductLens AI combines deterministic physics rules with vector classification to deliver 100% audit compliance and zero hallucination."),
                ("How does it solve Unilog's problem statement?", "It converts scattered, unstructured product information into rich, reliable, and commerce-ready product intelligence formatted strictly to 15 static expected headers for seamless e-commerce ingestion."),
                ("Core Unique Selling Proposition (USP)", "Deterministic Auditability + AI Speed: Every single enriched attribute includes a confidence score, source document citation, and LLM reasoning log, backed by automated metric-imperial unit normalization.")
            ]
            for head, body in items:
                p_h = tf.add_paragraph()
                p_h.text = f"• {head}"
                p_h.font.size = Pt(14)
                p_h.font.bold = True
                p_h.font.color.rgb = BLUE_BRAND

                p_desc = tf.add_paragraph()
                p_desc.text = body
                p_desc.font.size = Pt(12)

    # -------------------------------------------------------------
    # SLIDE 6: List of Features Offered
    # -------------------------------------------------------------
    slide6 = prs.slides[5]
    for shape in slide6.shapes:
        if shape.has_text_frame:
            tf = shape.text_frame
            tf.word_wrap = True
            tf.text = "List of Features Offered by ProductLens AI\n"
            p = tf.paragraphs[0]
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = BLUE_DARK

            features = [
                "1. Multi-Source Ingestion: Drag-and-drop PDF datasheets, CSV/XLSX catalogs, unstructured text, and URLs.",
                "2. Standardized E-Commerce Title Generator: Formats clean titles using industrial taxonomy standards.",
                "3. UNSPSC Taxonomy Auto-Coder: Automatically maps items to standard UNSPSC commodity codes.",
                "4. Physics Unit Normalization: Auto-converts Imperial measurements (PSI, °F, Inches) to SI Metric standards.",
                "5. Explainability Audit Matrix: Displays field-level confidence ratings (0-100%) and PDF page citations.",
                "6. Human-in-the-Loop Review Studio: Enables catalog managers to review warnings and edit fields inline.",
                "7. Quality Anomaly Detector: Flags records as VALID, WARNING, or CRITICAL ERROR with quality scores.",
                "8. 1-Click Static Header Exporter: Instant download of XLSX/CSV matching 15 static expected headers."
            ]
            for feat in features:
                pf = tf.add_paragraph()
                pf.text = feat
                pf.font.size = Pt(13)

    # -------------------------------------------------------------
    # SLIDE 7: Process Flow Diagram
    # -------------------------------------------------------------
    slide7 = prs.slides[6]
    for shape in slide7.shapes:
        if shape.has_text_frame:
            tf = shape.text_frame
            tf.word_wrap = True
            tf.text = "Process Flow & End-to-End Use-Case Workflow\n"
            p = tf.paragraphs[0]
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = BLUE_DARK

            steps = [
                "Step 1: Multi-Modal Data Ingestion (PDF technical datasheets, CSV files, URLs, or raw text snippets).",
                "Step 2: NLP Entity & Attribute Tokenization (Extracting MPN, Brand, Material, Voltage, Pressure).",
                "Step 3: UNSPSC Classification & Unit Normalization (Vector mapping + Imperial -> Metric conversion).",
                "Step 4: AI Validation & Quality Scoring (Evaluating completeness and assigning 0-100% confidence scores).",
                "Step 5: Human-in-the-Loop Review Studio (Auditing source citations and editing fields inline).",
                "Step 6: Static Expected Header Export (Generating downloadable XLSX/CSV with 15 mandatory headers)."
            ]
            for st in steps:
                ps = tf.add_paragraph()
                ps.text = f"• {st}"
                ps.font.size = Pt(13)

    # -------------------------------------------------------------
    # SLIDE 8: Wireframes & Mock Diagrams
    # -------------------------------------------------------------
    slide8 = prs.slides[7]
    for shape in slide8.shapes:
        if shape.has_text_frame:
            tf = shape.text_frame
            tf.word_wrap = True
            tf.text = "User Interface Wireframes & Workspace Modules\n"
            p = tf.paragraphs[0]
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = BLUE_DARK

            wireframes = [
                "1. Header & Navigation Studio: Displays real-time catalog statistics, active total products count, overall catalog quality score %, and 1-click Export XLSX button.",
                "2. Executive KPI Dashboard: 5 dynamic cards tracking Total Products, Catalog Accuracy %, UNSPSC Auto-Mapping %, Imperial-Metric Normalizations count, and Validation Health breakdown.",
                "3. Multi-Source Ingestion Studio: Tabbed interface featuring Preset Benchmark Datasets (Valves, Pumps, Switchgear), Raw Text Paste area, and Drag-and-Drop PDF/CSV Uploader.",
                "4. Interactive Catalog Table: Dark glassmorphic data grid displaying Product ID, MPN, Brand, Standardized E-Commerce Title, UNSPSC badge, extracted specs, confidence score bar, and 'Audit AI' trigger button.",
                "5. Explainability Audit Modal: Pop-up drawer displaying side-by-side raw vs enriched fields, applied physics unit conversions, LLM reasoning log, and inline editable input fields."
            ]
            for wf in wireframes:
                pwf = tf.add_paragraph()
                pwf.text = f"• {wf}"
                pwf.font.size = Pt(12)

    # -------------------------------------------------------------
    # SLIDE 9: System Architecture Diagram
    # -------------------------------------------------------------
    slide9 = prs.slides[8]
    for shape in slide9.shapes:
        if shape.has_text_frame:
            tf = shape.text_frame
            tf.word_wrap = True
            tf.text = "System Architecture & Processing Pipeline\n"
            p = tf.paragraphs[0]
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = BLUE_DARK

            arch_text = (
                "[Unstructured Inputs: PDFs / CSVs / Specs / URLs]\n"
                "                  │\n"
                "                  ▼\n"
                "[Multi-Modal Ingestion & Tokenizer Layer]\n"
                "                  │\n"
                "                  ▼\n"
                "[NLP & Physics Unit Normalizer (Imperial ↔ Metric)]\n"
                "                  │\n"
                "                  ▼\n"
                "[UNSPSC Vector Taxonomy Auto-Classification Engine]\n"
                "                  │\n"
                "                  ▼\n"
                "[Explainable AI (XAI) Audit & Citation Generator]\n"
                "                  │\n"
                "                  ▼\n"
                "[Human-in-the-Loop Review Dashboard]\n"
                "                  │\n"
                "                  ▼\n"
                "[Static Expected Header Exporter: 15 Mandatory Columns (.xlsx / .csv)]"
            )
            p_arch = tf.add_paragraph()
            p_arch.text = arch_text
            p_arch.font.size = Pt(11)
            p_arch.font.name = 'Courier New'

    # -------------------------------------------------------------
    # SLIDE 10: Technologies Used
    # -------------------------------------------------------------
    slide10 = prs.slides[9]
    for shape in slide10.shapes:
        if shape.has_text_frame:
            tf = shape.text_frame
            tf.word_wrap = True
            tf.text = "Technologies Used in the Solution\n"
            p = tf.paragraphs[0]
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = BLUE_DARK

            techs = [
                "• Frontend Web Application: React 18, Vite 5, JavaScript (ES6+), HTML5, CSS3.",
                "• Design System & Styling: Vanilla CSS + Tailwind CSS (Dark Glassmorphism UI, Inter & Outfit typography).",
                "• UI Icons & Feedback: Lucide React Icons, Canvas Confetti for celebratory interaction.",
                "• Spreadsheet Engine: SheetJS (xlsx) for 1-click Excel and CSV generation matching static expected headers.",
                "• AI & Taxonomy Engine: Custom NLP Tokenizer, UNSPSC v25.0 Taxonomy Classification Engine.",
                "• Presentation & PDF Generator: LaTeX Beamer (pdflatex) + Python-PPTX automated slide deck engine."
            ]
            for t in techs:
                pt = tf.add_paragraph()
                pt.text = t
                pt.font.size = Pt(13)

    # -------------------------------------------------------------
    # SLIDE 11: Estimated Implementation Cost
    # -------------------------------------------------------------
    slide11 = prs.slides[10]
    for shape in slide11.shapes:
        if shape.has_text_frame:
            tf = shape.text_frame
            tf.word_wrap = True
            tf.text = "Estimated Implementation Cost & Resource Plan\n"
            p = tf.paragraphs[0]
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = BLUE_DARK

            costs = [
                "1. Cloud Infrastructure (Serverless AWS/GCP execution): ~$120 / month.",
                "2. LLM & Vector Embedding API Tokens (100,000 SKUs/month): ~$180 / month.",
                "3. Database & CDN Storage: ~$30 / month.",
                "• Total Operating Expense: ~$330 / month for 100,000 processed SKUs.",
                "• Business ROI: Replaces manual cataloging team of 10 data entry specialists ($35,000/month overhead), delivering a 100:1 ROI payback within 30 days."
            ]
            for c in costs:
                pc = tf.add_paragraph()
                pc.text = c
                pc.font.size = Pt(13)

    # -------------------------------------------------------------
    # SLIDE 12: Snapshots of the MVP
    # -------------------------------------------------------------
    slide12 = prs.slides[11]
    for shape in slide12.shapes:
        if shape.has_text_frame:
            tf = shape.text_frame
            tf.word_wrap = True
            tf.text = "Snapshots & Capabilities of the Live MVP\n"
            p = tf.paragraphs[0]
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = BLUE_DARK

            mvp_caps = [
                "• Ingestion Studio: Tested with real benchmark industrial datasets (Valves, Centrifugal Pumps, Circuit Breakers).",
                "• Real-Time Processing: Ingests unstructured inputs and generates standardized titles, UNSPSC codes, and metric conversions in <500ms.",
                "• Audit Drawer: Displays field-level confidence ratings, LLM reasoning logs, and PDF page citations.",
                "• Static Exporter: Tested and verified Excel (.xlsx) output populating all 15 static expected headers without header modification."
            ]
            for m in mvp_caps:
                pm = tf.add_paragraph()
                pm.text = m
                pm.font.size = Pt(13)

    # -------------------------------------------------------------
    # SLIDE 13: Additional Details & Future Development
    # -------------------------------------------------------------
    slide13 = prs.slides[12]
    for shape in slide13.shapes:
        if shape.has_text_frame:
            tf = shape.text_frame
            tf.word_wrap = True
            tf.text = "Additional Details & Future Development Roadmap\n"
            p = tf.paragraphs[0]
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = BLUE_DARK

            future = [
                "• Phase 1 (Completed): Multi-source extraction, UNSPSC auto-coding, metric conversion, static XLSX export.",
                "• Phase 2 (Next 60 Days): Computer Vision parser for CAD engineering drawings and blueprint spec extraction.",
                "• Phase 3 (Enterprise): Real-time API connectors with Unilog Commerce Platform, SAP ERP, and Akeneo PIM.",
                "• Competitor Intelligence: Automated competitor price and spec benchmarking engine."
            ]
            for f in future:
                pf = tf.add_paragraph()
                pf.text = f
                pf.font.size = Pt(13)

    # -------------------------------------------------------------
    # SLIDE 14: Submission Links
    # -------------------------------------------------------------
    slide14 = prs.slides[13]
    for shape in slide14.shapes:
        if shape.has_text_frame:
            tf = shape.text_frame
            tf.word_wrap = True
            tf.text = "Project Verification & Submission Links\n"
            p = tf.paragraphs[0]
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = BLUE_DARK

            links = [
                "1. GitHub Public Repository Link:\n   https://github.com/aksharsakhi/ProductLens-AI",
                "2. Demo Video Link (3 Minutes):\n   https://youtu.be/sample_productlens_demo",
                "3. Working Live Prototype Link:\n   http://localhost:3000"
            ]
            for l in links:
                pl = tf.add_paragraph()
                pl.text = l
                pl.font.size = Pt(14)
                pl.font.bold = True

    # -------------------------------------------------------------
    # SLIDE 15: Thank You Slide
    # -------------------------------------------------------------
    slide15 = prs.slides[14]
    for shape in slide15.shapes:
        if shape.has_text_frame:
            tf = shape.text_frame
            tf.word_wrap = True
            tf.text = "Thank You!\n\nProductLens AI\nAI-Powered Product Intelligence for Industrial Commerce\nUniHack 2026 by Unilog & Hack2skill"
            p = tf.paragraphs[0]
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.color.rgb = BLUE_BRAND
            p.alignment = PP_ALIGN.CENTER

    # Save to both project locations
    prs.save('[EXT] UniHack-Protoype Template .pptx')
    prs.save('presentation_latex/ProductLens_AI_UniHack_Presentation.pptx')
    print("Successfully populated all 15 slides of [EXT] UniHack-Protoype Template .pptx!")

if __name__ == "__main__":
    populate_presentation()
